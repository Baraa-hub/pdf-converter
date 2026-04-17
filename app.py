from flask import Flask, request, send_file, render_template, jsonify
import os, uuid, zipfile, base64, re, unicodedata, traceback
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path
from PIL import Image
import pdfplumber

app = Flask(__name__, static_folder='static')
UPLOAD_FOLDER = '/tmp/uploads'
OUTPUT_FOLDER = '/tmp/outputs'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

def detect_pdf_type(input_path):
    try:
        with pdfplumber.open(input_path) as pdf:
            has_text = False
            has_images = False
            for page in pdf.pages:
                text = page.extract_text()
                if text and len(text.strip()) > 50:
                    has_text = True
                if page.images:
                    has_images = True
                if has_text and has_images:
                    break
            if has_text and has_images:
                return 'mixed'
            elif has_text:
                return 'text'
            else:
                return 'scanned'
    except:
        return 'scanned'

def pdf_to_images(input_path, dpi=120):
    return convert_from_path(input_path, dpi=dpi, thread_count=1, use_cropbox=True, strict=False)

def ocr_images(images):
    import pytesseract
    return [pytesseract.image_to_string(img, lang='eng+ara') for img in images]

def is_rtl_text(text):
    return any(unicodedata.bidirectional(c) in ('R', 'AL') for c in text if c.strip())

def fix_rtl(line):
    try:
        import arabic_reshaper
        line = arabic_reshaper.reshape(line)
    except:
        pass
    from bidi.algorithm import get_display
    return get_display(line)

def fix_arabic_for_docx(line):
    """For Word docs: reshape Arabic letters to connect them, but don't
    apply bidi display transform — Word handles RTL direction natively."""
    try:
        import arabic_reshaper
        line = arabic_reshaper.reshape(line)
    except:
        pass
    return line

def clean_text(text):
    if not text:
        return ''
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    text = text.replace('\n', ' ').replace('\r', ' ')
    text = re.sub(r' +', ' ', text).strip()
    return text

def save_image_file(img, path, save_fmt):
    img = img.copy()
    if save_fmt == 'JPEG':
        if img.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            if img.mode in ('RGBA', 'LA'):
                background.paste(img, mask=img.split()[-1])
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        img.save(path, 'JPEG', quality=95)
    else:
        if img.mode not in ('RGB', 'RGBA'):
            img = img.convert('RGB')
        img.save(path, 'PNG')

def save_as_docx_images(images, output_path, uid):
    from docx import Document
    from docx.shared import Inches, Pt, Emu
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    import lxml.etree as etree

    doc = Document()

    for i, img in enumerate(images):
        img_path = os.path.join(OUTPUT_FOLDER, f'{uid}_p{i+1}.png')
        save_image_file(img, img_path, 'PNG')
        w_px, h_px = img.size

        # Use current section for first page, add new section for subsequent pages
        if i == 0:
            section = doc.sections[0]
        else:
            # Add continuous section break then switch to new page section
            new_section = doc.add_section()
            section = doc.sections[-1]

        # Set page size exactly to image dimensions (at 96 DPI → EMU)
        # 1 inch = 914400 EMU, image rendered at ~100dpi
        dpi = 100
        page_w_emu = int(w_px / dpi * 914400)
        page_h_emu = int(h_px / dpi * 914400)
        section.page_width = Emu(page_w_emu)
        section.page_height = Emu(page_h_emu)
        section.top_margin = Emu(0)
        section.bottom_margin = Emu(0)
        section.left_margin = Emu(0)
        section.right_margin = Emu(0)

        # Add the image paragraph
        para = doc.add_paragraph()
        para.paragraph_format.space_before = Pt(0)
        para.paragraph_format.space_after = Pt(0)
        run = para.add_run()
        run.add_picture(img_path, width=Emu(page_w_emu), height=Emu(page_h_emu))

    doc.save(output_path)

def save_as_pptx_images(images, output_path, uid):
    from pptx import Presentation
    from pptx.util import Inches, Emu
    first_img = images[0]
    img_w, img_h = first_img.size
    slide_w = Inches(10)
    slide_h = Emu(int(slide_w * img_h / img_w))
    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    blank_layout = prs.slide_layouts[6]
    for i, img in enumerate(images):
        img_path = os.path.join(OUTPUT_FOLDER, f'{uid}_p{i+1}.png')
        save_image_file(img, img_path, 'PNG')
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, width=slide_w, height=slide_h)
    prs.save(output_path)

def save_as_html_images(images, output_path, uid):
    pages_html = ''
    for i, img in enumerate(images):
        img_path = os.path.join(OUTPUT_FOLDER, f'{uid}_p{i+1}.png')
        save_image_file(img, img_path, 'PNG')
        with open(img_path, 'rb') as f:
            b64 = base64.b64encode(f.read()).decode()
        pages_html += f'<div class="page"><img src="data:image/png;base64,{b64}" alt="Page {i+1}"></div>\n'
    html = f'''<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<style>
body{{margin:0;padding:20px;background:#666;display:flex;flex-direction:column;align-items:center}}
.page{{margin:10px 0;box-shadow:0 4px 12px rgba(0,0,0,0.4)}}
.page img{{display:block;max-width:900px;width:100%}}
</style></head>
<body>{pages_html}</body></html>'''
    with open(output_path, 'w') as f:
        f.write(html)

# ── Helpers ────────────────────────────────────────────────────────────────────

def rgb_to_hex(color):
    try:
        if isinstance(color, (list, tuple)) and len(color) == 3:
            r, g, b = [int(round(c * 255)) for c in color]
            return f'{r:02X}{g:02X}{b:02X}'
        if isinstance(color, (list, tuple)) and len(color) == 1:
            v = int(round(color[0] * 255))
            return f'{v:02X}{v:02X}{v:02X}'
    except:
        pass
    return None

def is_bold(word):
    fontname = word.get('fontname', '')
    return any(b in fontname.lower() for b in ['bold', 'black', 'heavy', 'semibold', 'demi'])

def apply_rtl_to_paragraph(para):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    pPr = para._p.get_or_add_pPr()
    bidi_el = OxmlElement('w:bidi')
    pPr.append(bidi_el)
    para.alignment = 2

def apply_rtl_to_run(run):
    from docx.oxml import OxmlElement
    rPr = run._r.get_or_add_rPr()
    rtl_el = OxmlElement('w:rtl')
    rPr.append(rtl_el)

def set_cell_background(cell, hex_color):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def compute_median_font_size(page):
    sizes = []
    try:
        words = page.extract_words(extra_attrs=['size'])
        for w in words:
            try:
                sizes.append(float(w.get('size', 12)))
            except:
                pass
    except:
        pass
    if not sizes:
        return 12.0
    sizes.sort()
    return sizes[len(sizes) // 2]

def get_rect_color_at(rects, y_top, y_bottom, x0, x1, page_w, page_h):
    best_color = None
    best_area = 0
    for r in rects:
        if r['width'] > page_w * 0.95 and r['height'] > page_h * 0.5:
            continue
        oy0 = max(r['top'], y_top)
        oy1 = min(r['bottom'], y_bottom)
        ox0 = max(r['x0'], x0)
        ox1 = min(r['x1'], x1)
        if oy1 > oy0 and ox1 > ox0:
            area = (oy1 - oy0) * (ox1 - ox0)
            if area > best_area:
                color = r.get('non_stroking_color')
                if color is not None:
                    hex_c = rgb_to_hex(color)
                    if hex_c and hex_c.upper() not in ('FFFFFF', 'FEFEFE', 'FDFDFD'):
                        best_color = hex_c
                        best_area = area
    return best_color

def merge_split_cells(row):
    """
    Merge cells where an English word is split mid-word across column boundary.
    Only merges when: cur ends with ASCII letter AND next starts with lowercase ASCII letter.
    Never merges when next cell starts with digit (digits = values, not continuations).
    """
    if not row:
        return row
    merged = list(row)
    i = 0
    while i < len(merged) - 1:
        cur = clean_text(str(merged[i]) if merged[i] else '')
        nxt = clean_text(str(merged[i + 1]) if merged[i + 1] else '')
        if not cur or not nxt:
            i += 1
            continue
        cur_ends_with_ascii_letter = cur[-1].isascii() and cur[-1].isalpha()
        nxt_starts_lowercase_ascii = nxt[0].isascii() and nxt[0].islower()
        cur_not_arabic = not is_rtl_text(cur)
        if cur_ends_with_ascii_letter and nxt_starts_lowercase_ascii and cur_not_arabic:
            merged[i] = cur + nxt
            merged.pop(i + 1)
        else:
            i += 1
    return merged

def extract_page_images_pymupdf(input_path, page_index, uid, output_folder):
    saved = []
    try:
        import fitz
        doc = fitz.open(input_path)
        page = doc[page_index]
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                img_bytes = base_image['image']
                img_ext = base_image.get('ext', 'png')
                w = base_image.get('width', 0)
                h = base_image.get('height', 0)
                if w < 50 or h < 50:
                    continue
                img_path = os.path.join(output_folder, f'{uid}_pg{page_index}_img{img_index}.{img_ext}')
                with open(img_path, 'wb') as f:
                    f.write(img_bytes)
                saved.append(img_path)
            except:
                pass
        doc.close()
    except ImportError:
        pass
    except Exception:
        pass
    return saved

# ── Debug endpoint ─────────────────────────────────────────────────────────────

@app.route('/debug', methods=['POST'])
def debug():
    """Upload a PDF and get a diagnostic report of what the server sees."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    uid = str(uuid.uuid4())[:8]
    input_path = os.path.join(UPLOAD_FOLDER, f'{uid}.pdf')
    file.save(input_path)

    report = {'pages': []}
    try:
        with pdfplumber.open(input_path) as pdf:
            report['page_count'] = len(pdf.pages)
            for page_index, page in enumerate(pdf.pages[:2]):  # Only check first 2 pages
                page_w = float(page.width)
                page_h = float(page.height)
                rects = page.rects

                cell_rects = [r for r in rects if
                              r['width'] < page_w * 0.95 and
                              r['height'] < page_h * 0.3 and
                              r['width'] > 5 and r['height'] > 5]

                page_info = {
                    'page': page_index + 1,
                    'total_rects': len(rects),
                    'cell_rects': len(cell_rects),
                    'tables_found': 0,
                    'table_rows': 0,
                    'table_cols': 0,
                    'error': None,
                    'sample_row': None,
                }

                if cell_rects:
                    y_pos = sorted(set(
                        [round(r['top'], 1) for r in cell_rects] +
                        [round(r['bottom'], 1) for r in cell_rects]
                    ))
                    x_pos = sorted(set(
                        [round(r['x0'], 1) for r in cell_rects] +
                        [round(r['x1'], 1) for r in cell_rects]
                    ))
                    page_info['y_pos_count'] = len(y_pos)
                    page_info['x_pos_count'] = len(x_pos)

                    try:
                        tables = page.extract_tables({
                            'vertical_strategy': 'explicit',
                            'horizontal_strategy': 'explicit',
                            'explicit_vertical_lines': x_pos,
                            'explicit_horizontal_lines': y_pos,
                            'snap_tolerance': 4,
                            'join_tolerance': 4,
                        }) or []
                        page_info['tables_found'] = len(tables)
                        if tables:
                            page_info['table_rows'] = len(tables[0])
                            page_info['table_cols'] = len(tables[0][0]) if tables[0] else 0
                            if tables[0]:
                                page_info['sample_row'] = [str(c)[:20] if c else '' for c in tables[0][2]] if len(tables[0]) > 2 else []
                    except Exception as e:
                        page_info['error'] = traceback.format_exc()

                report['pages'].append(page_info)
    except Exception as e:
        report['fatal_error'] = traceback.format_exc()
    finally:
        if os.path.exists(input_path):
            os.remove(input_path)

    return jsonify(report)

# ── Main native DOCX converter ─────────────────────────────────────────────────

def save_as_docx_native(input_path, output_path, uid, pages_param=None):
    from docx import Document
    from docx.shared import Pt, Inches

    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)

    with pdfplumber.open(input_path) as pdf:
        total_pages = len(pdf.pages)
        page_indices = parse_pages(pages_param, total_pages) if pages_param else list(range(total_pages))
        first_page = True
        for page_index in page_indices:
            page = pdf.pages[page_index]
            if not first_page:
                doc.add_page_break()
            first_page = False

            page_w = float(page.width)
            page_h = float(page.height)
            rects = page.rects
            median_size = compute_median_font_size(page)

            # Embed page images — but skip full-page images (they cause blank pages)
            img_paths = extract_page_images_pymupdf(input_path, page_index, uid, OUTPUT_FOLDER)
            for img_path in img_paths:
                try:
                    pil_img = Image.open(img_path)
                    w_px, h_px = pil_img.size
                    if w_px < 50 or h_px < 50:
                        continue
                    # Skip images that are close to full-page size — these are
                    # page backgrounds/renders, not embedded photos
                    page_aspect = page_w / page_h if page_h > 0 else 1
                    img_aspect = w_px / h_px if h_px > 0 else 1
                    aspect_similar = abs(page_aspect - img_aspect) < 0.1
                    # If image is large AND same aspect ratio as page → skip
                    if aspect_similar and w_px > 400 and h_px > 400:
                        continue
                    doc.add_picture(img_path, width=Inches(6.5))
                except:
                    pass

            # Detect tables via filled rects
            cell_rects = [r for r in rects if
                          r['width'] < page_w * 0.95 and
                          r['height'] < page_h * 0.3 and
                          r['width'] > 5 and r['height'] > 5]

            has_table = False
            tables = []
            y_pos = []
            x_pos = []

            if cell_rects:
                y_pos = sorted(set(
                    [round(r['top'], 1) for r in cell_rects] +
                    [round(r['bottom'], 1) for r in cell_rects]
                ))
                x_pos = sorted(set(
                    [round(r['x0'], 1) for r in cell_rects] +
                    [round(r['x1'], 1) for r in cell_rects]
                ))
                if len(y_pos) >= 2 and len(x_pos) >= 2:
                    try:
                        tables = page.extract_tables({
                            'vertical_strategy': 'explicit',
                            'horizontal_strategy': 'explicit',
                            'explicit_vertical_lines': x_pos,
                            'explicit_horizontal_lines': y_pos,
                            'snap_tolerance': 4,
                            'join_tolerance': 4,
                        }) or []
                        if tables:
                            has_table = True
                    except:
                        pass

            if not has_table:
                try:
                    tables = page.extract_tables({
                        'vertical_strategy': 'lines',
                        'horizontal_strategy': 'lines',
                        'snap_tolerance': 3,
                    }) or []
                    if tables:
                        has_table = True
                        y_pos = []
                        x_pos = []
                except:
                    pass

            if has_table and tables:
                table_data = tables[0]
                rows = [r for r in table_data if any(c and str(c).strip() for c in r)]

                if rows:
                    rows = [merge_split_cells(r) for r in rows]
                    num_cols = max(len(r) for r in rows)
                    norm_rows = [list(r) + [None] * (num_cols - len(r)) for r in rows]

                    tbl = doc.add_table(rows=len(norm_rows), cols=num_cols)
                    tbl.style = 'Table Grid'

                    row_y_spans = []
                    if len(y_pos) >= 2:
                        row_y_spans = [(y_pos[i], y_pos[i+1]) for i in range(len(y_pos)-1)]

                    col_x_spans = []
                    if len(x_pos) >= 2:
                        col_x_spans = [(x_pos[j], x_pos[j+1]) for j in range(min(num_cols, len(x_pos)-1))]

                    orig_indices = [i for i, r in enumerate(table_data)
                                    if any(c and str(c).strip() for c in r)]

                    for r_idx, row in enumerate(norm_rows):
                        orig_r = orig_indices[r_idx] if r_idx < len(orig_indices) else r_idx
                        y_top, y_bottom = row_y_spans[orig_r] if orig_r < len(row_y_spans) else (0, 0)

                        for c_idx, cell_val in enumerate(row):
                            cell = tbl.rows[r_idx].cells[c_idx]
                            text = clean_text(str(cell_val) if cell_val else '')
                            if text.lower() == 'none':
                                text = ''

                            if y_top != y_bottom and c_idx < len(col_x_spans):
                                cx0, cx1 = col_x_spans[c_idx]
                                hex_color = get_rect_color_at(
                                    rects, y_top, y_bottom, cx0, cx1, page_w, page_h)
                                if hex_color:
                                    try:
                                        set_cell_background(cell, hex_color)
                                    except:
                                        pass

                            if not text:
                                continue

                            rtl = is_rtl_text(text)
                            if rtl:
                                text = fix_rtl(text)

                            para = cell.paragraphs[0]
                            if rtl:
                                apply_rtl_to_paragraph(para)
                            run = para.add_run(text)
                            run.font.size = Pt(10)
                            if rtl:
                                apply_rtl_to_run(run)

                    doc.add_paragraph()

            else:
                try:
                    all_words = page.extract_words(
                        x_tolerance=3, y_tolerance=3,
                        keep_blank_chars=False,
                        use_text_flow=False,
                        extra_attrs=['fontname', 'size']
                    )
                except:
                    all_words = []

                if not all_words:
                    try:
                        import pytesseract
                        imgs = convert_from_path(input_path, dpi=150,
                                                first_page=page_index+1,
                                                last_page=page_index+1)
                        if imgs:
                            text = pytesseract.image_to_string(imgs[0], lang='eng+ara')
                            for line in text.split('\n'):
                                line = clean_text(line)
                                if line:
                                    para = doc.add_paragraph()
                                    rtl = is_rtl_text(line)
                                    if rtl:
                                        line = fix_rtl(line)
                                        apply_rtl_to_paragraph(para)
                                    run = para.add_run(line)
                                    run.font.size = Pt(11)
                                    if rtl:
                                        apply_rtl_to_run(run)
                    except:
                        pass
                    continue

                lines_dict = {}
                for word in all_words:
                    y_key = round(float(word['top']) / 4) * 4
                    if y_key not in lines_dict:
                        lines_dict[y_key] = []
                    lines_dict[y_key].append(word)

                for y_key in sorted(lines_dict.keys()):
                    words = sorted(lines_dict[y_key], key=lambda w: float(w['x0']))
                    line_text = clean_text(' '.join(w['text'] for w in words))
                    if not line_text:
                        continue

                    sizes = []
                    for w in words:
                        try:
                            sizes.append(float(w.get('size', 12)))
                        except:
                            sizes.append(12.0)
                    avg_size = sum(sizes) / len(sizes) if sizes else 12.0
                    bold_count = sum(1 for w in words if is_bold(w))
                    mostly_bold = bold_count > len(words) / 2
                    rtl = is_rtl_text(line_text)

                    if rtl:
                        line_text = fix_rtl(line_text)

                    para = doc.add_paragraph()
                    para.paragraph_format.space_before = Pt(0)
                    para.paragraph_format.space_after = Pt(2)
                    if rtl:
                        apply_rtl_to_paragraph(para)

                    run = para.add_run(line_text)
                    run.bold = mostly_bold or (avg_size > median_size * 1.3)
                    run.font.size = Pt(min(max(8, avg_size), 72))
                    if rtl:
                        apply_rtl_to_run(run)

    doc.save(output_path)


def save_as_docx_text(input_path, output_path, pages_param=None):
    """OCR-based DOCX — always renders pages to images then runs tesseract.
    This avoids garbled text from PDFs with custom font encodings."""
    from docx import Document
    from docx.shared import Pt, Inches
    import pytesseract

    doc = Document()

    with pdfplumber.open(input_path) as pdf:
        total_pages = len(pdf.pages)
        page_indices = parse_pages(pages_param, total_pages) if pages_param else list(range(total_pages))
        first_page = True

        for i in page_indices:
            page = pdf.pages[i]
            if not first_page:
                doc.add_page_break()
            first_page = False

            section = doc.sections[-1]
            section.page_width = Inches(float(page.width) / 72)
            section.page_height = Inches(float(page.height) / 72)
            section.top_margin = Inches(0.75)
            section.bottom_margin = Inches(0.75)
            section.left_margin = Inches(0.75)
            section.right_margin = Inches(0.75)

            # Always render to image and OCR — avoids garbled text from
            # PDFs with custom/private-use font encodings
            imgs = convert_from_path(input_path, dpi=200,
                                     first_page=i+1, last_page=i+1,
                                     use_cropbox=True, strict=False)
            if imgs:
                text = pytesseract.image_to_string(imgs[0], lang='eng+ara')
                for line in text.split('\n'):
                    line = clean_text(line)
                    if line:
                        para = doc.add_paragraph()
                        rtl = is_rtl_text(line)
                        if rtl:
                            line = fix_arabic_for_docx(line)
                            apply_rtl_to_paragraph(para)
                        run = para.add_run(line)
                        run.font.size = Pt(11)
                        if rtl:
                            apply_rtl_to_run(run)

    doc.save(output_path)


def save_as_xlsx(input_path, output_path, pages_param, page_count):
    """
    Scan PDF for tables and export each table to a separate Excel sheet.
    Uses rect-based detection (same as DOCX native) for accuracy.
    Raises ValueError if no tables found.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    wb.remove(wb.active)  # Remove default empty sheet
    sheet_count = 0

    with pdfplumber.open(input_path) as pdf:
        total_pages = len(pdf.pages)
        page_indices = parse_pages(pages_param, total_pages) if pages_param else list(range(total_pages))

        for page_index in page_indices:
            page = pdf.pages[page_index]
            page_w = float(page.width)
            page_h = float(page.height)
            rects = page.rects

            # Try rect-based table extraction first
            cell_rects = [r for r in rects if
                          r['width'] < page_w * 0.95 and
                          r['height'] < page_h * 0.3 and
                          r['width'] > 5 and r['height'] > 5]

            tables = []
            y_pos = []
            x_pos = []

            if cell_rects:
                y_pos = sorted(set(
                    [round(r['top'], 1) for r in cell_rects] +
                    [round(r['bottom'], 1) for r in cell_rects]
                ))
                x_pos = sorted(set(
                    [round(r['x0'], 1) for r in cell_rects] +
                    [round(r['x1'], 1) for r in cell_rects]
                ))
                if len(y_pos) >= 2 and len(x_pos) >= 2:
                    try:
                        tables = page.extract_tables({
                            'vertical_strategy': 'explicit',
                            'horizontal_strategy': 'explicit',
                            'explicit_vertical_lines': x_pos,
                            'explicit_horizontal_lines': y_pos,
                            'snap_tolerance': 4,
                            'join_tolerance': 4,
                        }) or []
                    except:
                        pass

            # Fallback to line-based
            if not tables:
                try:
                    tables = page.extract_tables({
                        'vertical_strategy': 'lines',
                        'horizontal_strategy': 'lines',
                        'snap_tolerance': 3,
                    }) or []
                    y_pos = []
                    x_pos = []
                except:
                    pass

            for t_idx, table_data in enumerate(tables):
                # Filter empty rows
                rows = [r for r in table_data if any(c and str(c).strip() for c in r)]
                if not rows:
                    continue

                sheet_count += 1
                sheet_name = f'P{page_index+1}_T{t_idx+1}'
                ws = wb.create_sheet(title=sheet_name)

                # Get row y spans for color detection
                row_y_spans = []
                if len(y_pos) >= 2:
                    row_y_spans = [(y_pos[i], y_pos[i+1]) for i in range(len(y_pos)-1)]
                col_x_spans = []
                if len(x_pos) >= 2:
                    col_x_spans = [(x_pos[j], x_pos[j+1]) for j in range(len(x_pos)-1)]

                orig_indices = [i for i, r in enumerate(table_data)
                                if any(c and str(c).strip() for c in r)]

                # Styles
                thin_border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )

                for r_idx, row in enumerate(rows):
                    orig_r = orig_indices[r_idx] if r_idx < len(orig_indices) else r_idx
                    y_top, y_bottom = row_y_spans[orig_r] if orig_r < len(row_y_spans) else (0, 0)

                    for c_idx, cell_val in enumerate(row):
                        text = clean_text(str(cell_val) if cell_val else '')
                        if text.lower() == 'none':
                            text = ''

                        xl_cell = ws.cell(row=r_idx+1, column=c_idx+1, value=text)
                        xl_cell.border = thin_border
                        xl_cell.alignment = Alignment(wrap_text=True, vertical='center')

                        # Detect background color from rects
                        if y_top != y_bottom and c_idx < len(col_x_spans):
                            cx0, cx1 = col_x_spans[c_idx]
                            hex_color = get_rect_color_at(
                                rects, y_top, y_bottom, cx0, cx1, page_w, page_h)
                            if hex_color:
                                try:
                                    xl_cell.fill = PatternFill(
                                        start_color=hex_color,
                                        end_color=hex_color,
                                        fill_type='solid'
                                    )
                                    # White text on dark backgrounds
                                    r_val = int(hex_color[0:2], 16)
                                    g_val = int(hex_color[2:4], 16)
                                    b_val = int(hex_color[4:6], 16)
                                    brightness = (r_val * 299 + g_val * 587 + b_val * 114) / 1000
                                    if brightness < 128:
                                        xl_cell.font = Font(color='FFFFFF', bold=True)
                                    else:
                                        xl_cell.font = Font(bold=True)
                                except:
                                    pass

                # Auto-fit column widths
                for col in ws.columns:
                    max_len = 0
                    col_letter = get_column_letter(col[0].column)
                    for cell in col:
                        try:
                            if cell.value:
                                max_len = max(max_len, len(str(cell.value)))
                        except:
                            pass
                    ws.column_dimensions[col_letter].width = min(max(max_len + 2, 8), 50)

    if sheet_count == 0:
        raise ValueError('No tables detected in this PDF. Try a PDF that contains structured tables with visible borders.')

    wb.save(output_path)


def get_page_images(input_path, pages_param, page_count, dpi):
    """Get PIL images for either all pages or specific pages."""
    if pages_param:
        indices = parse_pages(pages_param, page_count)
        result = []
        for idx in indices:
            imgs = convert_from_path(
                input_path, dpi=dpi, thread_count=1,
                first_page=idx+1, last_page=idx+1,
                use_cropbox=True, strict=False
            )
            if imgs:
                result.append(imgs[0])
        return result
    else:
        return pdf_to_images(input_path, dpi=dpi)

def parse_pages(pages_param, total):
    indices = []
    for part in pages_param.split(','):
        part = part.strip()
        if '-' in part:
            try:
                parts = part.split('-')
                start = int(parts[0].strip())
                end = int(parts[1].strip())
                indices += list(range(start - 1, end))
            except:
                pass
        else:
            try:
                indices.append(int(part) - 1)
            except:
                pass
    return sorted(set([i for i in indices if 0 <= i < total]))

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/privacy')
def privacy():
    return render_template('privacy.html')

# ── Merge PDFs ──────────────────────────────────────────────────────────────
@app.route('/merge', methods=['POST'])
def merge_pdfs():
    from pypdf import PdfWriter
    uid = str(uuid.uuid4())[:8]
    saved_paths = []
    output_path = os.path.join(OUTPUT_FOLDER, f'{uid}_merged.pdf')
    try:
        files = request.files.getlist('files')
        if not files or len(files) < 2:
            return jsonify({'error': 'Please upload at least 2 PDF files'}), 400
        writer = PdfWriter()
        for f in files:
            if not f.filename.lower().endswith('.pdf'):
                return jsonify({'error': f'{f.filename} is not a PDF'}), 400
            p = os.path.join(UPLOAD_FOLDER, f'{uid}_{len(saved_paths)}.pdf')
            f.save(p)
            saved_paths.append(p)
            from pypdf import PdfReader
            reader = PdfReader(p)
            for page in reader.pages:
                writer.add_page(page)
        with open(output_path, 'wb') as out:
            writer.write(out)
        return send_file(output_path, as_attachment=True, download_name='merged.pdf')
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        for p in saved_paths:
            try:
                if os.path.exists(p): os.remove(p)
            except: pass
        try:
            if os.path.exists(output_path): os.remove(output_path)
        except: pass

# ── Split PDF ───────────────────────────────────────────────────────────────
@app.route('/split', methods=['POST'])
def split_pdf():
    from pypdf import PdfReader, PdfWriter
    uid = str(uuid.uuid4())[:8]
    input_path = os.path.join(UPLOAD_FOLDER, f'{uid}.pdf')
    output_zip = os.path.join(OUTPUT_FOLDER, f'{uid}_split.zip')
    output_paths = []
    try:
        f = request.files.get('file')
        if not f:
            return jsonify({'error': 'No file uploaded'}), 400
        split_type = request.form.get('split_type', 'all')  # 'all' or 'range'
        f.save(input_path)
        reader = PdfReader(input_path)
        total = len(reader.pages)

        if split_type == 'all':
            # Every page becomes its own PDF
            ranges = [(i, i) for i in range(total)]
        else:
            # Parse range like "1-3, 5-7, 9"
            pages_param = request.form.get('pages', '').strip()
            ranges = []
            for part in pages_param.split(','):
                part = part.strip()
                if '-' in part:
                    s, e = part.split('-')
                    ranges.append((int(s.strip())-1, int(e.strip())-1))
                elif part:
                    n = int(part) - 1
                    ranges.append((n, n))

        if not ranges:
            return jsonify({'error': 'No valid page ranges specified'}), 400

        for idx, (start, end) in enumerate(ranges):
            writer = PdfWriter()
            for page_num in range(max(0, start), min(end+1, total)):
                writer.add_page(reader.pages[page_num])
            out_path = os.path.join(OUTPUT_FOLDER, f'{uid}_part{idx+1}.pdf')
            with open(out_path, 'wb') as out:
                writer.write(out)
            output_paths.append((out_path, f'part_{idx+1}_pages_{start+1}-{end+1}.pdf'))

        if len(output_paths) == 1:
            # Single output — return directly
            return send_file(output_paths[0][0], as_attachment=True,
                           download_name=output_paths[0][1])
        else:
            # Multiple outputs — zip them
            with zipfile.ZipFile(output_zip, 'w', zipfile.ZIP_DEFLATED) as zf:
                for path, arcname in output_paths:
                    zf.write(path, arcname)
            return send_file(output_zip, as_attachment=True, download_name='split_pages.zip')
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        try:
            if os.path.exists(input_path): os.remove(input_path)
        except: pass
        for path, _ in output_paths:
            try:
                if os.path.exists(path): os.remove(path)
            except: pass
        try:
            if os.path.exists(output_zip): os.remove(output_zip)
        except: pass

# ── Compress PDF ────────────────────────────────────────────────────────────
@app.route('/compress', methods=['POST'])
def compress_pdf():
    import subprocess
    uid = str(uuid.uuid4())[:8]
    input_path = os.path.join(UPLOAD_FOLDER, f'{uid}.pdf')
    output_path = os.path.join(OUTPUT_FOLDER, f'{uid}_compressed.pdf')
    try:
        f = request.files.get('file')
        if not f:
            return jsonify({'error': 'No file uploaded'}), 400
        level = request.form.get('level', 'medium')  # low, medium, high
        f.save(input_path)
        original_size = os.path.getsize(input_path)

        # Use Ghostscript if available, otherwise pypdf
        gs_settings = {
            'low':    '/printer',
            'medium': '/ebook',
            'high':   '/screen',
        }
        gs_setting = gs_settings.get(level, '/ebook')

        gs_result = subprocess.run([
            'gs', '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4',
            f'-dPDFSETTINGS={gs_setting}',
            '-dNOPAUSE', '-dQUIET', '-dBATCH',
            f'-sOutputFile={output_path}', input_path
        ], capture_output=True, timeout=120)

        if gs_result.returncode != 0 or not os.path.exists(output_path):
            # Fallback: pypdf compression
            from pypdf import PdfReader, PdfWriter
            reader = PdfReader(input_path)
            writer = PdfWriter()
            for page in reader.pages:
                page.compress_content_streams()
                writer.add_page(page)
            with open(output_path, 'wb') as out:
                writer.write(out)

        compressed_size = os.path.getsize(output_path)
        reduction = round((1 - compressed_size / original_size) * 100, 1)

        response = send_file(output_path, as_attachment=True,
                           download_name='compressed.pdf')
        response.headers['X-Original-Size'] = str(original_size)
        response.headers['X-Compressed-Size'] = str(compressed_size)
        response.headers['X-Reduction'] = str(reduction)
        return response
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        try:
            if os.path.exists(input_path): os.remove(input_path)
        except: pass
        try:
            if os.path.exists(output_path): os.remove(output_path)
        except: pass


def convert_office_to_pdf(input_path, output_dir):
    """Convert DOCX/XLSX/PPTX to PDF using LibreOffice headless."""
    import subprocess
    result = subprocess.run([
        'libreoffice', '--headless', '--convert-to', 'pdf',
        '--outdir', output_dir, input_path
    ], capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise Exception(f'LibreOffice conversion failed: {result.stderr}')
    # LibreOffice names output as original_filename.pdf
    base = os.path.splitext(os.path.basename(input_path))[0]
    out = os.path.join(output_dir, base + '.pdf')
    if not os.path.exists(out):
        raise Exception('Output PDF not found after conversion')
    return out

def images_to_pdf(image_paths, output_path):
    """Combine multiple images into a single PDF."""
    images = []
    for p in image_paths:
        img = Image.open(p).convert('RGB')
        images.append(img)
    if not images:
        raise Exception('No valid images found')
    images[0].save(output_path, save_all=True, append_images=images[1:])

@app.route('/convert-to-pdf', methods=['POST'])
def convert_to_pdf():
    fmt = request.form.get('format', '').lower()
    uid = str(uuid.uuid4())[:8]
    saved_paths = []

    try:
        files = request.files.getlist('files')
        if not files or all(f.filename == '' for f in files):
            return jsonify({'error': 'No files uploaded'}), 400

        if fmt == 'image':
            # Multiple images → single PDF
            allowed = {'.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp', '.tiff'}
            img_paths = []
            for i, f in enumerate(files):
                ext = os.path.splitext(f.filename.lower())[1]
                if ext not in allowed:
                    return jsonify({'error': f'Unsupported image format: {ext}'}), 400
                p = os.path.join(UPLOAD_FOLDER, f'{uid}_{i}{ext}')
                f.save(p)
                saved_paths.append(p)
                img_paths.append(p)

            base_name = secure_filename(files[0].filename)
            base_name = os.path.splitext(base_name)[0].replace('.', '_')
            output_filename = f'{base_name}_converted.pdf'
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            images_to_pdf(img_paths, output_path)

        elif fmt in ('docx', 'xlsx', 'pptx'):
            # Single office file → PDF
            if len(files) != 1:
                return jsonify({'error': 'Please upload exactly one file'}), 400
            f = files[0]
            ext = os.path.splitext(f.filename.lower())[1]
            expected = {'docx': '.docx', 'xlsx': '.xlsx', 'pptx': '.pptx'}
            if ext != expected[fmt]:
                return jsonify({'error': f'Expected a {expected[fmt]} file'}), 400

            input_path = os.path.join(UPLOAD_FOLDER, f'{uid}{ext}')
            f.save(input_path)
            saved_paths.append(input_path)

            base_name = secure_filename(f.filename)
            base_name = os.path.splitext(base_name)[0].replace('.', '_')
            output_filename = f'{base_name}_converted.pdf'

            pdf_path = convert_office_to_pdf(input_path, OUTPUT_FOLDER)
            # Rename to our output filename
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            if pdf_path != output_path:
                os.rename(pdf_path, output_path)

        else:
            return jsonify({'error': 'Unsupported format'}), 400

        return send_file(output_path, as_attachment=True, download_name=output_filename)

    except Exception as e:
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500
    finally:
        for p in saved_paths:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except:
                pass
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except:
            pass

@app.route('/detect', methods=['POST'])
def detect():
    if 'file' not in request.files:
        return jsonify({'error': 'No file'}), 400
    file = request.files['file']
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'Not a PDF'}), 400
    uid = str(uuid.uuid4())[:8]
    input_path = os.path.join(UPLOAD_FOLDER, f'{uid}.pdf')
    file.save(input_path)
    pdf_type = detect_pdf_type(input_path)
    with pdfplumber.open(input_path) as pdf:
        page_count = len(pdf.pages)
    return jsonify({'type': pdf_type, 'uid': uid, 'page_count': page_count})

@app.route('/preview-pptx', methods=['POST'])
def preview_pptx():
    """Convert a specific slide of PPTX to image using LibreOffice."""
    import subprocess
    if 'file' not in request.files:
        return jsonify({'error': 'No file'}), 400
    f = request.files['file']
    slide_num = int(request.form.get('slide', 1))
    uid = str(uuid.uuid4())[:8]
    input_path = os.path.join(UPLOAD_FOLDER, f'{uid}.pptx')
    f.save(input_path)
    try:
        # Convert PPTX to PDF using LibreOffice
        result = subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'pdf',
            '--outdir', OUTPUT_FOLDER, input_path
        ], capture_output=True, text=True, timeout=60)
        lo_pdf = os.path.join(OUTPUT_FOLDER, f'{uid}.pdf')
        if not os.path.exists(lo_pdf):
            return jsonify({'error': 'Conversion failed'}), 500
        # Get total slide count
        with pdfplumber.open(lo_pdf) as pdf:
            total_slides = len(pdf.pages)
        # Render requested slide as image
        imgs = convert_from_path(lo_pdf, dpi=150,
                                 first_page=slide_num, last_page=slide_num)
        if not imgs:
            return jsonify({'error': 'Could not render slide'}), 500
        img_path = os.path.join(OUTPUT_FOLDER, f'{uid}_slide.jpg')
        imgs[0].save(img_path, 'JPEG', quality=85)
        # Return image with total slides in header
        response = send_file(img_path, mimetype='image/jpeg')
        response.headers['X-Total-Slides'] = str(total_slides)
        return response
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        for p in [input_path,
                  os.path.join(OUTPUT_FOLDER, f'{uid}.pdf'),
                  os.path.join(OUTPUT_FOLDER, f'{uid}_slide.jpg')]:
            try:
                if os.path.exists(p): os.remove(p)
            except: pass


@app.route('/convert', methods=['POST'])
def convert():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    file = request.files['file']
    fmt = request.form.get('format', 'jpg').lower()
    mode = request.form.get('mode', 'image')
    if file.filename == '' or not file.filename.endswith('.pdf'):
        return jsonify({'error': 'Please upload a valid PDF'}), 400

    uid = str(uuid.uuid4())[:8]
    input_path = os.path.join(UPLOAD_FOLDER, f'{uid}.pdf')
    file.save(input_path)
    base_name = secure_filename(file.filename).replace('.pdf', '').replace('.', '_')

    try:
        with pdfplumber.open(input_path) as pdf:
            page_count = len(pdf.pages)
        dpi = 150 if page_count <= 5 else 100

        if fmt in ('jpg', 'png'):
            save_fmt = 'JPEG' if fmt == 'jpg' else 'PNG'
            ext = fmt
            pages_param = request.form.get('pages', '').strip()
            if pages_param:
                # Parse first, then only render the specific requested pages
                selected_indices = parse_pages(pages_param, page_count)
                if not selected_indices:
                    return jsonify({'error': 'No valid pages selected. Use format like: 1, 3, 5-7'}), 400
                selected_images = []
                for idx in selected_indices:
                    imgs = convert_from_path(
                        input_path, dpi=dpi, thread_count=1,
                        first_page=idx+1, last_page=idx+1,
                        use_cropbox=True, strict=False
                    )
                    if imgs:
                        selected_images.append(imgs[0])
            else:
                # All pages
                images = pdf_to_images(input_path, dpi=dpi)
                selected_images = images

            if len(selected_images) == 1:
                output_filename = f'{base_name}_converted.{ext}'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                save_image_file(selected_images[0], output_path, save_fmt)
            else:
                output_filename = f'{base_name}_converted.zip'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                img_paths = []
                for i, img in enumerate(selected_images):
                    img_path = os.path.join(OUTPUT_FOLDER, f'{uid}_p{i+1}.{ext}')
                    save_image_file(img, img_path, save_fmt)
                    img_paths.append((img_path, f'page{i+1}.{ext}'))
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                    for img_path, arcname in img_paths:
                        zf.write(img_path, arcname)

        elif fmt == 'docx':
            output_filename = f'{base_name}_converted.docx'
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            pages_param = request.form.get('pages', '').strip()
            if mode == 'ocr':
                save_as_docx_text(input_path, output_path, pages_param or None)
            elif mode == 'image':
                page_images = get_page_images(input_path, pages_param, page_count, dpi)
                save_as_docx_images(page_images, output_path, uid)
            else:
                save_as_docx_native(input_path, output_path, uid, pages_param or None)

        elif fmt == 'pptx':
            output_filename = f'{base_name}_converted.pptx'
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            pages_param = request.form.get('pages', '').strip()
            page_images = get_page_images(input_path, pages_param, page_count, dpi)
            if mode == 'ocr':
                from pptx import Presentation
                from pptx.util import Inches, Emu
                pages_text = ocr_images(page_images)
                first_img = page_images[0]
                slide_w = Inches(10)
                slide_h = Emu(int(slide_w * first_img.size[1] / first_img.size[0]))
                prs = Presentation()
                prs.slide_width = slide_w
                prs.slide_height = slide_h
                slide_layout = prs.slide_layouts[1]
                for i, text in enumerate(pages_text):
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = f'Page {i+1}'
                    tf = slide.placeholders[1].text_frame
                    tf.word_wrap = True
                    tf.text = text[:800] if text.strip() else '(no text detected)'
                prs.save(output_path)
            else:
                save_as_pptx_images(page_images, output_path, uid)

        elif fmt == 'html':
            output_filename = f'{base_name}_converted.html'
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            pages_param = request.form.get('pages', '').strip()
            page_images = get_page_images(input_path, pages_param, page_count, dpi)
            save_as_html_images(page_images, output_path, uid)

        elif fmt == 'xlsx':
            output_filename = f'{base_name}_converted.xlsx'
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            pages_param = request.form.get('pages', '').strip()
            save_as_xlsx(input_path, output_path, pages_param or None, page_count)

        else:
            return jsonify({'error': 'Unsupported format'}), 400

        return send_file(output_path, as_attachment=True, download_name=output_filename)

    except Exception as e:
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500
    finally:
        # Delete uploaded PDF immediately
        if os.path.exists(input_path):
            os.remove(input_path)
        # Delete output file and any temp page images after sending
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except:
            pass
        # Clean up any temp page images created during conversion
        import glob
        for tmp_file in glob.glob(os.path.join(OUTPUT_FOLDER, f'{uid}_*')):
            try:
                os.remove(tmp_file)
            except:
                pass

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port)
